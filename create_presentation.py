from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def add_title_with_style(slide, title_text):
    title = slide.shapes.title
    title.text = title_text
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(40)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(31, 73, 125)

def add_code_snippet(slide, code_text, left=Inches(0.5), top=Inches(2.5), width=Inches(9), height=Inches(4)):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = code_text
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = 'Consolas'
        paragraph.font.size = Pt(14)
    
    # Add a semi-transparent rectangle behind the code
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
    shape.line.color.rgb = RGBColor(200, 200, 200)
    shape.zorder = 0
    textbox.zorder = 1

def create_presentation():
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_title_with_style(slide, "Deep Learning Interactive Visualization")
    subtitle = slide.placeholders[1]
    subtitle.text = "Martin Dyrba\nGerman Center for Neurodegenerative Diseases (DZNE)"
    
    # Add app screenshot to title slide
    if os.path.exists("InteractiveVis.png"):
        left = Inches(1)
        top = Inches(4)
        pic = slide.shapes.add_picture("InteractiveVis.png", left, top, height=Inches(3))
    
    # Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_with_style(slide, "Project Overview")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "• Purpose: CNN model for Alzheimer's disease detection\n"
    tf.add_paragraph().text = "• Interactive visualization of brain regions relevance\n"
    tf.add_paragraph().text = "• Published in Alzheimer's Research & Therapy (2021)\n"
    tf.add_paragraph().text = "• DOI: 10.1186/s13195-021-00924-2"
    
    # Key Features with Code
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_with_style(slide, "Key Features & Implementation")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "• 3D Convolutional Neural Network\n"
    tf.add_paragraph().text = "• Interactive Brain Region Visualization\n"
    tf.add_paragraph().text = "• Real-time Relevance Mapping"
    
    code_snippet = """def apply_thresholds(self, relevance_map, threshold=0.5, cluster_size=20):
    # Apply threshold to relevance map
    self.overlay = np.copy(relevance_map)
    self.overlay[np.abs(self.overlay) < threshold] = 0
    
    # Cluster size filtering
    labelimg = np.copy(self.overlay)
    labelimg[labelimg > 0] = 1  # binarize img
    labelimg = label(labelimg, connectivity=2)
    
    # Calculate cluster properties
    lprops = regionprops(labelimg, intensity_image=self.overlay)
    self.clust_sizes = []
    for lab in lprops:
        if lab.area < cluster_size:
            labelimg[labelimg == lab.label] = 0  # remove small clusters"""
    
    add_code_snippet(slide, code_snippet)
    
    # Visualization Implementation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_with_style(slide, "Visualization Implementation")
    
    code_snippet = """def overlay2rgba(relevance_map, alpha=0.5):
    # Convert relevance map to RGBA visualization
    alpha_mask = np.copy(relevance_map)
    alpha_mask[np.abs(alpha_mask) > 0] = alpha
    
    # Scale to color range
    relevance_map = relevance_map / 2 + 0.5
    ovl = np.uint8(overlay_colormap(relevance_map) * 255)
    ovl[:, :, :, 3] = np.uint8(alpha_mask * 255)
    
    return ovl.view("uint32").reshape(ovl.shape[:3])"""
    
    add_code_snippet(slide, code_snippet)
    
    # Interactive Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_with_style(slide, "Interactive Features")
    
    code_snippet = """def click_frontal_callback(event):
    # Handle user interaction with frontal brain view
    if event.x < 1:
        x = 1
    elif event.x > slice_slider_sagittal.end:
        x = slice_slider_sagittal.end
    else:
        x = int(round(event.x))
        
    # Update other views
    slice_slider_sagittal.update(value=x)
    slice_slider_axial.update(value=y)
    plot_sagittal()"""
    
    add_code_snippet(slide, code_snippet)
    
    # Technical Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_with_style(slide, "Technical Architecture")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "• Model-View-Controller (MVC) Architecture\n"
    tf.add_paragraph().text = "• Bokeh Web Application Framework\n"
    tf.add_paragraph().text = "• TensorFlow for CNN Implementation\n"
    tf.add_paragraph().text = "• Docker Containerization"
    
    # Implementation Details
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_with_style(slide, "Implementation Details")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "• Python-based implementation\n"
    tf.add_paragraph().text = "• TensorFlow 1.15 framework\n"
    tf.add_paragraph().text = "• Bokeh web application\n"
    tf.add_paragraph().text = "• Docker containerization\n"
    tf.add_paragraph().text = "• GPU-accelerated training support"
    
    # Usage & Deployment
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_with_style(slide, "Usage & Deployment")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "• Public web service (explaination.net/demo)\n"
    tf.add_paragraph().text = "• Docker container deployment\n"
    tf.add_paragraph().text = "• Local installation\n"
    tf.add_paragraph().text = "• Multiple deployment options for different use cases"
    
    if os.path.exists("InteractiveVis.png"):
        left = Inches(6)
        top = Inches(2)
        pic = slide.shapes.add_picture("InteractiveVis.png", left, top, height=Inches(4))
    
    # Results & Impact
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_with_style(slide, "Results & Impact")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "• Improved CNN comprehensibility\n"
    tf.add_paragraph().text = "• Interactive relevance map visualization\n"
    tf.add_paragraph().text = "• Validation across multiple datasets\n"
    tf.add_paragraph().text = "• Clinical application potential\n"
    tf.add_paragraph().text = "• Enhanced understanding of model decisions"
    
    # Resources & Links
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_with_style(slide, "Resources & Links")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "• GitHub: github.com/martindyrba/DeepLearningInteractiveVis\n"
    tf.add_paragraph().text = "• Docker Hub: docker pull martindyrba/interactivevis\n"
    tf.add_paragraph().text = "• Demo: explaination.net/demo\n"
    tf.add_paragraph().text = "• Documentation available in project README"
    
    # Save the presentation
    prs.save('Project_Presentation_New.pptx')

if __name__ == "__main__":
    create_presentation() 